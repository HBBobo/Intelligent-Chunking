"""
Extract text from Office documents (PPTX, DOCX).
"""

from pathlib import Path


def normalize_text_encoding(text: str) -> str:
    """
    Normalize text by replacing problematic Windows-1252 characters
    with their proper UTF-8 equivalents.

    Args:
        text: Text that may contain Windows-1252 characters.

    Returns:
        Text with normalized encoding.
    """
    # Common Windows-1252 to UTF-8 mappings
    replacements = {
        '\x96': '–',  # en-dash
        '\x97': '—',  # em-dash
        '\x91': ''',  # left single quote
        '\x92': ''',  # right single quote
        '\x93': '"',  # left double quote
        '\x94': '"',  # right double quote
        '\x85': '…',  # ellipsis
        '\x95': '•',  # bullet
        '\xa0': ' ',  # non-breaking space
        '\xad': '-',  # soft hyphen
    }

    for old, new in replacements.items():
        text = text.replace(old, new)

    # Also try to encode/decode to clean up any remaining issues
    try:
        # Encode to UTF-8 with replacement for any remaining bad chars
        text = text.encode('utf-8', errors='replace').decode('utf-8')
    except Exception:
        pass

    return text


def extract_pptx_to_text(pptx_path: Path) -> str:
    """
    Extract text from a PowerPoint file.

    Args:
        pptx_path: Path to the .pptx file.

    Returns:
        Extracted text with slide separations.
    """
    try:
        from pptx import Presentation
    except ImportError:
        raise ImportError(
            "python-pptx not installed. Install with: pip install python-pptx"
        )

    prs = Presentation(pptx_path)
    slides_text = []

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_content = []

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text = shape.text.strip()
                # Check if this looks like a title (first text, short)
                if not slide_content and len(text) < 100:
                    slide_content.append(text)  # Title as heading
                else:
                    slide_content.append(text)

        if slide_content:
            slides_text.append("\n".join(slide_content))

    text = "\n\n".join(slides_text)
    return normalize_text_encoding(text)


def extract_docx_to_text(docx_path: Path) -> str:
    """
    Extract text from a Word document.

    Args:
        docx_path: Path to the .docx file.

    Returns:
        Extracted text with paragraph separations.
    """
    try:
        from docx import Document
    except ImportError:
        raise ImportError(
            "python-docx not installed. Install with: pip install python-docx"
        )

    doc = Document(docx_path)
    paragraphs = []

    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            # Check if paragraph has heading style
            style_name = para.style.name.lower() if para.style else ""
            if "heading" in style_name or "title" in style_name:
                paragraphs.append(text)  # Will be detected as heading by segmenter
            else:
                paragraphs.append(text)

    text = "\n\n".join(paragraphs)
    return normalize_text_encoding(text)
